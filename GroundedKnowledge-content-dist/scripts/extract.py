"""Per-file text extraction + folder-taxonomy metadata for the Seismic corpus.

Importable by ingest.py. Text-only: no OCR, no image rasterization.
Each extractor yields "segments" — (locator_kind, locator, segment_type, text) —
which ingest.py chunks and embeds.
"""
import os
import re
import unicodedata
import subprocess
import tempfile
import json

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# --- taxonomy vocab (matched case-insensitively against path segments) ---
PERSONAS = [
    "Office of the CDO", "Cloud Modernization", "Cloud Only IT",
    "ERP_SAP Modernization", "ERP SAP Modernization",
]
_KNOWN_CONTENT_TYPES = {
    "demo video", "learning material", "presentation", "reference architecture",
    "operations material", "customer profile", "executive brief", "article",
    "ebook", "infographic", "solution brief", "product", "data sheet",
    "assessment tool", "customer slide", "video",
}


def rel_path(full_path, corpus_root):
    try:
        return os.path.relpath(full_path, corpus_root)
    except ValueError:
        return full_path


def parse_taxonomy(full_path, corpus_root):
    """Derive structured metadata from the folder path (the only metadata we have)."""
    rp = rel_path(full_path, corpus_root)
    parts = [p for p in rp.split(os.sep) if p and p not in (".", "..")]
    crumbs = parts[:-1]  # drop filename
    low = [c.lower() for c in crumbs]

    division = content_type = persona = use_case_kit = None
    audience = "Unknown"

    # Pricing Center pack (quote-gated licensing/pricing docs synced from Drive):
    # tag the whole subtree Internal and derive content_type from the subfolder.
    for i, c in enumerate(low):
        if "pricing center" in c:
            division = "Pricing Center"
            audience = "Internal"
            if i + 1 < len(crumbs):
                content_type = crumbs[i + 1]
            return {
                "division": division, "content_type": content_type or "Pricing Center",
                "persona": None, "use_case_kit": None, "audience": audience,
                "folder_crumbs": crumbs,
            }

    # division = first segment under "All Sales"
    if "all sales" in low:
        i = low.index("all sales")
        if i + 1 < len(crumbs):
            division = crumbs[i + 1]

    # audience split present under "Content by Type"
    for c in low:
        if c == "internal":
            audience = "Internal"
        elif c == "external":
            audience = "External"

    # persona
    for c in crumbs:
        for p in PERSONAS:
            if c.lower() == p.lower():
                persona = "ERP_SAP Modernization" if "erp" in p.lower() else p
    # use-case kit (segment after "On Demand Prospecting Kits")
    for i, c in enumerate(low):
        if "prospecting kits" in c and i + 1 < len(crumbs):
            use_case_kit = crumbs[i + 1]
            break

    # content_type = a known type appearing anywhere in the crumbs (prefer the deepest)
    for c in reversed(crumbs):
        if c.lower() in _KNOWN_CONTENT_TYPES:
            content_type = c
            break
    if content_type is None and division:
        content_type = division

    return {
        "division": division,
        "content_type": content_type,
        "persona": persona,
        "use_case_kit": use_case_kit,
        "audience": audience,
        "folder_crumbs": crumbs,
    }


def detect_language(text):
    """Cheap heuristic: CJK density for JA, stopword ratio for en/de/it. No new deps."""
    if not text:
        return "unknown"
    sample = text[:1200]
    cjk = sum(1 for ch in sample if "CJK" in unicodedata.name(ch, "") or
              "HIRAGANA" in unicodedata.name(ch, "") or "KATAKANA" in unicodedata.name(ch, ""))
    if cjk > len(sample) * 0.15:
        return "ja"
    words = re.findall(r"[a-zà-ÿäöüß]+", sample.lower())
    if not words:
        return "unknown"
    ws = set(words)
    de = {"und", "der", "die", "das", "für", "mit", "von", "ist", "den", "auf", "eine", "im"}
    it = {"il", "di", "che", "per", "con", "una", "sono", "gli", "nel", "come", "della"}
    en = {"the", "and", "for", "with", "that", "this", "from", "your", "how", "are", "our"}
    score = {"de": len(ws & de), "it": len(ws & it), "en": len(ws & en)}
    best = max(score, key=score.get)
    return best if score[best] > 0 else "en"


# --- per-type extractors: each yields dicts describing one segment ---
def _pdf_segments(path):
    doc = fitz.open(path)
    try:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            has_images = bool(page.get_images())
            needs_ocr = len(text) < 40 and has_images
            if text:
                yield {"page_number": i, "segment": "body", "text": text, "needs_ocr": needs_ocr}
            elif needs_ocr:
                yield {"page_number": i, "segment": "body", "text": "", "needs_ocr": True}
    finally:
        doc.close()


def _pptx_segments(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        body = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                body.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        body.append(" | ".join(cells))
        if body:
            yield {"slide_number": i, "segment": "body", "text": "\n".join(body)}
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                yield {"slide_number": i, "segment": "notes", "text": notes}


def _docx_segments(path):
    """Extract from .docx files. Falls back to plain text if not a real Word doc."""
    try:
        doc = DocxDocument(path)
        buf, section = [], 1
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            if para.style and para.style.name and para.style.name.lower().startswith("heading") and buf:
                yield {"page_number": section, "segment": "body", "text": "\n".join(buf)}
                buf, section = [], section + 1
            buf.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    buf.append(" | ".join(cells))
        if buf:
            yield {"page_number": section, "segment": "body", "text": "\n".join(buf)}
    except Exception as e:
        # Some .doc files are actually plain text - try reading as text
        if "not found" in str(e).lower() or "package" in str(e).lower():
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read().strip()
                    if text:
                        yield {"page_number": 1, "segment": "body", "text": text}
            except:
                raise  # Re-raise original error if text fallback also fails
        else:
            raise


def _xlsx_segments(path):
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    for si, ws in enumerate(wb.worksheets, start=1):
        lines = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                lines.append("\t".join(vals))
        if lines:
            yield {"page_number": si, "segment": "sheet", "text": f"[{ws.title}]\n" + "\n".join(lines)}
    wb.close()


def _video_segments(path):
    """Transcribe video using Whisper (tiny model for speed).

    Whisper requires ffmpeg for video/audio processing. If ffmpeg is not installed,
    videos are indexed with metadata only (filename, path) but no transcript.
    """
    try:
        import whisper
    except ImportError:
        # If Whisper not installed, skip gracefully with filename info
        basename = os.path.basename(path)
        yield {"page_number": 1, "segment": "video",
               "text": f"Video: {basename} [Whisper not installed - install with: pip install openai-whisper]"}
        return

    try:
        # Use tiny model for speed (base/small/medium available for better accuracy)
        model = whisper.load_model("tiny")

        # Transcribe with error handling for missing ffmpeg
        result = model.transcribe(path, language="en", fp16=False)

        # Yield full transcript
        full_text = result.get("text", "").strip()
        if full_text:
            yield {"page_number": 1, "segment": "transcript", "text": full_text}

        # Yield timestamped segments for better localization
        segments = result.get("segments", [])
        for i, seg in enumerate(segments, start=1):
            seg_text = seg.get("text", "").strip()
            if seg_text:
                yield {
                    "page_number": i,
                    "segment": "transcript_segment",
                    "text": seg_text,
                    "ts_start": seg.get("start"),
                    "ts_end": seg.get("end")
                }
    except FileNotFoundError as e:
        # ffmpeg not found - videos indexed by filename only
        basename = os.path.basename(path)
        if "ffmpeg" in str(e).lower():
            yield {"page_number": 1, "segment": "video",
                   "text": f"Video: {basename} [Transcription requires ffmpeg - install with: brew install ffmpeg]"}
        else:
            yield {"page_number": 1, "segment": "video", "text": f"Video: {basename} [Error: {str(e)}]"}
    except Exception as e:
        # Other errors - log but don't fail the whole index
        basename = os.path.basename(path)
        yield {"page_number": 1, "segment": "video", "text": f"Video: {basename} [Error: {str(e)[:100]}]"}


_EXTRACTORS = {
    ".pdf": _pdf_segments, ".pptx": _pptx_segments,
    ".docx": _docx_segments, ".doc": _docx_segments,
    ".xlsx": _xlsx_segments, ".xlsm": _xlsx_segments,
    ".mp4": _video_segments, ".mov": _video_segments,
    ".avi": _video_segments, ".mkv": _video_segments, ".webm": _video_segments,
}
SUPPORTED_EXTS = set(_EXTRACTORS)


def extract(path):
    """Return list of segment dicts for a supported file, or [] if unsupported/unreadable."""
    ext = os.path.splitext(path)[1].lower()
    fn = _EXTRACTORS.get(ext)
    if not fn:
        return []
    return list(fn(path))
